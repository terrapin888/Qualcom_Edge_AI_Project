# services/attachment_service.py - ì²¨ë¶€íŒŒì¼ ì²˜ë¦¬ ì„œë¹„ìŠ¤

import os
import io
import tempfile
import hashlib
from pathlib import Path
import numpy as np

#0824 ìˆ˜ì •
from services.genie_qwen import genie_summarize_document
#0824 ë

# ì„ íƒì  ì„í¬íŠ¸ - ì—†ëŠ” ë¼ì´ë¸ŒëŸ¬ë¦¬ëŠ” ë¹„í™œì„±í™”
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("[âš ï¸ PIL/Pillow ì—†ìŒ - ì´ë¯¸ì§€ ì²˜ë¦¬ ë¹„í™œì„±í™”]")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    print("[âš ï¸ pdfplumber ì—†ìŒ - PDF ì²˜ë¦¬ ë¹„í™œì„±í™”]")

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    print("[âš ï¸ PyPDF2 ì—†ìŒ - PDF ë°±ì—… ì²˜ë¦¬ ë¹„í™œì„±í™”]")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("[âš ï¸ python-docx ì—†ìŒ - Word ë¬¸ì„œ ì²˜ë¦¬ ë¹„í™œì„±í™”]")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("[âš ï¸ python-pptx ì—†ìŒ - PowerPoint ì²˜ë¦¬ ë¹„í™œì„±í™”]")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    print("[âš ï¸ pandas ì—†ìŒ - Excel ì²˜ë¦¬ ë¹„í™œì„±í™”]")

try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("[âš ï¸ pdf2image ì—†ìŒ - PDF OCR ì²˜ë¦¬ ë¹„í™œì„±í™”]")

class AttachmentService:
    def __init__(self, config, ai_models):
        self.config = config
        self.ai_models = ai_models
        self.attachment_cache = {}
        
        # ì‚¬ìš© ê°€ëŠ¥í•œ ê¸°ëŠ¥ ì²´í¬
        self.features = {
            'image_processing': PIL_AVAILABLE,
            'pdf_processing': PDFPLUMBER_AVAILABLE or PYPDF2_AVAILABLE,
            'docx_processing': DOCX_AVAILABLE,
            'pptx_processing': PPTX_AVAILABLE,
            'xlsx_processing': PANDAS_AVAILABLE,
            'pdf_ocr': PDF2IMAGE_AVAILABLE,
            'yolo': hasattr(ai_models, 'load_yolo_model'),
            'ocr': hasattr(ai_models, 'load_ocr_model')
        }
        
        print(f"[ğŸ“ ì²¨ë¶€íŒŒì¼ ì„œë¹„ìŠ¤ ì´ˆê¸°í™”] ì‚¬ìš© ê°€ëŠ¥í•œ ê¸°ëŠ¥: {sum(self.features.values())}/{len(self.features)}")
    
    def process_email_attachments(self, email_message, email_subject, email_id):
        """ì´ë©”ì¼ì—ì„œ ì²¨ë¶€íŒŒì¼ì„ ì¶”ì¶œí•˜ê³  ì²˜ë¦¬ (ìºì‹± í¬í•¨)"""
        cache_key = f"email_{email_id}"
        
        # ìºì‹œ í™•ì¸
        if cache_key in self.attachment_cache:
            print(f"[ğŸ“ ìºì‹œ ì‚¬ìš©] {email_subject[:30]}...")
            return self.attachment_cache[cache_key]
        
        attachments = []
        print(f"[ğŸ“ ìƒˆë¡œìš´ ì²¨ë¶€íŒŒì¼ ì²˜ë¦¬] {email_subject[:30]}...")
        
        try:
            for part in email_message.walk():
                if part.get_content_disposition() == 'attachment':
                    attachment_info = self._process_single_attachment(part, email_subject)
                    if attachment_info:
                        attachments.append(attachment_info)
        except Exception as e:
            print(f"[â—ì²¨ë¶€íŒŒì¼ ì›Œí‚¹ ì˜¤ë¥˜] {str(e)}")
        
        # ìºì‹œ ì €ì¥
        self.attachment_cache[cache_key] = attachments
        self._manage_cache_size()
        
        print(f"[âœ… ì²¨ë¶€íŒŒì¼ ì²˜ë¦¬ ì™„ë£Œ] {len(attachments)}ê°œ ì²˜ë¦¬ë¨")
        return attachments
    
    def _process_single_attachment(self, part, email_subject):
        """ê°œë³„ ì²¨ë¶€íŒŒì¼ ì²˜ë¦¬"""
        try:
            filename = self._decode_filename(part.get_filename())
            if not filename:
                return None
            
            attachment_data = part.get_payload(decode=True)
            if not attachment_data:
                return None
            
            file_ext = Path(filename).suffix.lower()
            mime_type = part.get_content_type()
            
            attachment_info = {
                'filename': filename,
                'size': len(attachment_data),
                'mime_type': mime_type,
                'extension': file_ext
            }
            
            # íŒŒì¼ íƒ€ì…ë³„ ì²˜ë¦¬
            if file_ext in {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}:
                attachment_info.update(self._process_image(attachment_data, filename))
            elif file_ext == '.pdf' or 'pdf' in mime_type:
                attachment_info.update(self._process_pdf(attachment_data, filename))
            elif file_ext == '.docx' or 'wordprocessingml' in mime_type:
                attachment_info.update(self._process_docx(attachment_data, filename))
            elif file_ext == '.pptx' or 'presentationml' in mime_type:
                attachment_info.update(self._process_pptx(attachment_data, filename))
            elif file_ext in ['.xlsx', '.xls'] or 'spreadsheetml' in mime_type:
                attachment_info.update(self._process_xlsx(attachment_data, filename))
            else:
                attachment_info.update({'type': 'other', 'processing_method': 'metadata_only'})
            
            return attachment_info
            
        except Exception as e:
            print(f"[â—ì²¨ë¶€íŒŒì¼ ì²˜ë¦¬ ì˜¤ë¥˜] {filename if 'filename' in locals() else 'Unknown'}: {str(e)}")
            return None
    
    def _decode_filename(self, filename):
        """íŒŒì¼ëª… ë””ì½”ë”©"""
        if not filename:
            return None
        
        try:
            from email.header import decode_header
            decoded_parts = decode_header(filename)
            if decoded_parts and decoded_parts[0]:
                decoded_filename = decoded_parts[0]
                if isinstance(decoded_filename[0], bytes):
                    return decoded_filename[0].decode(decoded_filename[1] or 'utf-8')
                else:
                    return decoded_filename[0]
        except:
            pass
        
        return filename
    
    def _process_image(self, attachment_data, filename):
        """ì´ë¯¸ì§€ ì²˜ë¦¬ (YOLO + OCR)"""
        try:
            if not PIL_AVAILABLE:
                return {'type': 'image', 'error': 'PIL not available', 'processing_method': 'disabled'}
            
            # YOLO ê°ì²´ ì¸ì‹ (ONNX ìš°ì„ , PyTorch í´ë°±)
            yolo_detections = []
            if self.features['yolo']:
                # ONNX YOLO ì‹œë„
                if hasattr(self.ai_models, 'yolo_onnx_session') and self.ai_models.yolo_onnx_session:
                    print("[ğŸš€ YOLO ONNX] ê°ì²´ íƒì§€ ì‹œì‘...")
                    yolo_detections = self._yolo_detect_objects_onnx(attachment_data)
                # PyTorch YOLO í´ë°±
                elif self.ai_models.load_yolo_model():
                    print("[ğŸš€ YOLO PyTorch] ê°ì²´ íƒì§€ ì‹œì‘ (í´ë°±)...")
                    yolo_detections = self._yolo_detect_objects(attachment_data)
            
            # OCR í…ìŠ¤íŠ¸ ì¶”ì¶œ
            ocr_result = {'text': '', 'success': False}
            print(f"[ğŸ” OCR ì²´í¬] features['ocr']: {self.features['ocr']}")
            if self.features['ocr'] and self.ai_models.load_ocr_model():
                print(f"[ğŸ” OCR ì‹œì‘] {filename}")
                ocr_result = self._extract_text_with_ocr(attachment_data, filename)
                print(f"[ğŸ” OCR ê²°ê³¼] success: {ocr_result.get('success')}, text length: {len(ocr_result.get('text', ''))}")
            else:
                print(f"[âš ï¸ OCR ê±´ë„ˆëœ€] features['ocr']: {self.features['ocr']}")
            
            result = {
                'type': 'image',
                'yolo_detections': yolo_detections,
                'detected_objects': [det['class'] for det in yolo_detections],
                'object_count': len(yolo_detections),
                'extracted_text': ocr_result.get('text', ''),
                'ocr_success': ocr_result.get('success', False),
                'processing_method': f"YOLO({len(yolo_detections)}) + OCR({ocr_result.get('success', False)})"
            }
            
            # í…ìŠ¤íŠ¸ ìš”ì•½ ìƒì„±
            if ocr_result.get('success') and ocr_result.get('text'):
                result['text_summary'] = self._summarize_document(
                    ocr_result['text'], filename, 'image_with_text'
                )
            
            return result
            
        except Exception as e:
            print(f"[â—ì´ë¯¸ì§€ ì²˜ë¦¬ ì˜¤ë¥˜] {str(e)}")
            return {'type': 'image', 'error': str(e), 'processing_method': 'failed'}
    
    def _yolo_detect_objects(self, image_data):
        """YOLO ê°ì²´ ì¸ì‹"""
        try:
            if not PIL_AVAILABLE:
                return []
            
            # ì´ë¯¸ì§€ ë¡œë“œ ë° ì „ì²˜ë¦¬
            image = Image.open(io.BytesIO(image_data))
            
            # RGBA â†’ RGB ë³€í™˜
            if image.mode in ['RGBA', 'LA']:
                rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                image = rgb_image
            elif image.mode != 'RGB':
                image = image.convert('RGB')
            
            image_np = np.array(image)
            
            # YOLO ì¶”ë¡ 
            results = self.ai_models.yolo_model(image_np, conf=0.5)
            
            detections = []
            if len(results) > 0 and results[0].boxes is not None:
                boxes = results[0].boxes
                for i in range(len(boxes)):
                    conf = float(boxes.conf[i].cpu().numpy())
                    cls = int(boxes.cls[i].cpu().numpy())
                    class_name = self.ai_models.yolo_model.names[cls]
                    
                    detections.append({
                        'class': class_name,
                        'confidence': conf,
                        'class_id': cls
                    })
            
            return detections
            
        except Exception as e:
            print(f"[â—YOLO ì²˜ë¦¬ ì˜¤ë¥˜] {str(e)}")
            return []
    
    def _yolo_detect_objects_onnx(self, image_data):
        """YOLO ONNX ê°ì²´ ì¸ì‹"""
        try:
            if not PIL_AVAILABLE:
                return []
            
            # ì´ë¯¸ì§€ ë¡œë“œ ë° ì „ì²˜ë¦¬
            image = Image.open(io.BytesIO(image_data))
            
            # RGBA â†’ RGB ë³€í™˜
            if image.mode in ['RGBA', 'LA']:
                rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                image = rgb_image
            elif image.mode != 'RGB':
                image = image.convert('RGB')
            
            image_np = np.array(image)
            
            # YOLO ONNX ì¶”ë¡ 
            detections = self.ai_models.detect_objects_with_yolo_onnx(image_np)
            
            # bbox ì •ë³´ ì œê±°í•˜ê³  attachment_serviceì—ì„œ ì‚¬ìš©í•˜ëŠ” í˜•ì‹ìœ¼ë¡œ ë³€í™˜
            simplified_detections = []
            for det in detections:
                simplified_detections.append({
                    'class': det['class'],
                    'confidence': det['confidence'],
                    'class_id': det['class_id']
                })
            
            return simplified_detections
            
        except Exception as e:
            print(f"[â—YOLO ONNX ì²˜ë¦¬ ì˜¤ë¥˜] {str(e)}")
            return []
    
    def _extract_text_with_ocr(self, attachment_data, filename):
        """OCR í…ìŠ¤íŠ¸ ì¶”ì¶œ"""
        try:
            if not PIL_AVAILABLE:
                print(f"[â—OCR] PIL ì‚¬ìš© ë¶ˆê°€ëŠ¥")
                return {'text': '', 'success': False, 'error': 'PIL not available'}
            
            print(f"[ğŸ” OCR ì‹œì‘] íŒŒì¼ëª…: {filename}, ë°ì´í„° í¬ê¸°: {len(attachment_data)} bytes")
            
            image = Image.open(io.BytesIO(attachment_data))
            print(f"[ğŸ” OCR] ì›ë³¸ ì´ë¯¸ì§€ ëª¨ë“œ: {image.mode}, í¬ê¸°: {image.size}")
            
            # ì´ë¯¸ì§€ ì „ì²˜ë¦¬
            if image.mode in ['RGBA', 'LA']:
                print(f"[ğŸ” OCR] ì´ë¯¸ì§€ ëª¨ë“œ ë³€í™˜: {image.mode} -> RGB")
                rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                image = rgb_image
            elif image.mode != 'RGB':
                print(f"[ğŸ” OCR] ì´ë¯¸ì§€ ëª¨ë“œ ë³€í™˜: {image.mode} -> RGB")
                image = image.convert('RGB')
            
            print(f"[ğŸ” OCR] ìµœì¢… ì´ë¯¸ì§€ ëª¨ë“œ: {image.mode}, í¬ê¸°: {image.size}")
            image_np = np.array(image)
            print(f"[ğŸ” OCR] numpy ë°°ì—´ í˜•íƒœ: {image_np.shape}, dtype: {image_np.dtype}")
            
            # OCR ìˆ˜í–‰ (ONNX ìš°ì„ , EasyOCR API í´ë°±)
            result = None
            
            # NPU ìš°ì„ , ONNX í´ë°± ì²˜ë¦¬ (extract_text_from_image_onnxê°€ ë‚´ë¶€ì ìœ¼ë¡œ NPU ìš°ì„  ì²˜ë¦¬í•¨)
            if ((hasattr(self.ai_models, 'npu_detector_session') and self.ai_models.npu_detector_session and
                 hasattr(self.ai_models, 'npu_recognizer_session') and self.ai_models.npu_recognizer_session) or
                (hasattr(self.ai_models, 'easyocr_detector_session') and self.ai_models.easyocr_detector_session and
                 hasattr(self.ai_models, 'easyocr_recognizer_session') and self.ai_models.easyocr_recognizer_session)):
                print(f"[ğŸš€ OCR] NPU/ONNX í…ìŠ¤íŠ¸ ì¶”ì¶œ ì‹œì‘...")
                result = self.ai_models.extract_text_from_image_onnx(image_np)
                if result:
                    print(f"[ğŸ“‹ Attachment] OCR ê²°ê³¼ ìˆ˜ì‹  - {len(result)}ê°œ í…ìŠ¤íŠ¸ ì˜ì—­")
                    if result:
                        print(f"[ğŸ“Š Attachment] ì²« ë²ˆì§¸ ê²°ê³¼: {result[0]}")
                else:
                    print(f"[âŒ Attachment] OCR ê²°ê³¼ ì—†ìŒ")
                    result = []
            else:
                detector_loaded = hasattr(self.ai_models, 'easyocr_detector_session') and self.ai_models.easyocr_detector_session
                recognizer_loaded = hasattr(self.ai_models, 'easyocr_recognizer_session') and self.ai_models.easyocr_recognizer_session
                print(f"[âŒ OCR] ONNX OCR ëª¨ë¸ì´ ë¡œë”©ë˜ì§€ ì•ŠìŒ - Detector: {detector_loaded}, Recognizer: {recognizer_loaded}")
                result = []
            
            if result:
                print(f"[ğŸ” OCR] ì²« ë²ˆì§¸ ê²°ê³¼ ì˜ˆì‹œ: {result[0] if result else 'None'}")
            
            text = ""
            confident_detections = 0
            for i, detection in enumerate(result):
                print(f"[ğŸ” OCR] íƒì§€ {i+1}: {detection}")
                if len(detection) >= 3:
                    text_content = detection[1]
                    confidence = detection[2]
                    print(f"[ğŸ” OCR] í…ìŠ¤íŠ¸: '{text_content}', ì‹ ë¢°ë„: {confidence:.3f}")
                    if confidence > 0.3:
                        text += text_content + " "
                        confident_detections += 1
                        print(f"[âœ… OCR] ì¶”ê°€ë¨ (ì‹ ë¢°ë„ {confidence:.3f}): {text_content}")
                    else:
                        print(f"[âŒ OCR] ë‚®ì€ ì‹ ë¢°ë„ë¡œ ì œì™¸: {text_content}")
            
            final_text = text.strip()
            success = bool(final_text)
            
            print(f"[ğŸ” OCR ìµœì¢… ê²°ê³¼] ì„±ê³µ: {success}, ì‹ ë¢°ë„ ë†’ì€ íƒì§€: {confident_detections}, ìµœì¢… í…ìŠ¤íŠ¸ ê¸¸ì´: {len(final_text)}")
            print(f"[ğŸ” OCR ìµœì¢… í…ìŠ¤íŠ¸] '{final_text[:100]}{'...' if len(final_text) > 100 else ''}'")
            
            return {
                'text': final_text,
                'success': success,
                'method': 'ocr',
                'total_detections': len(result),
                'confident_detections': confident_detections
            }
            
        except Exception as e:
            print(f"[â—OCR ì˜ˆì™¸] {str(e)}")
            import traceback
            print(f"[â—OCR ìŠ¤íƒíŠ¸ë ˆì´ìŠ¤] {traceback.format_exc()}")
            return {'text': '', 'success': False, 'error': str(e)}
    
    def _process_pdf(self, attachment_data, filename):
        """PDF ì²˜ë¦¬"""
        if not PDFPLUMBER_AVAILABLE and not PYPDF2_AVAILABLE:
            return {'type': 'document_pdf', 'error': 'PDF libraries not available', 'extraction_success': False}
        
        try:
            # pdfplumberë¡œ í…ìŠ¤íŠ¸ ì¶”ì¶œ ì‹œë„
            if PDFPLUMBER_AVAILABLE:
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
                    temp_file.write(attachment_data)
                    temp_file_path = temp_file.name
                
                try:
                    with pdfplumber.open(temp_file_path) as pdf:
                        text = ""
                        for page_num, page in enumerate(pdf.pages):
                            page_text = page.extract_text()
                            if page_text:
                                text += f"\n=== í˜ì´ì§€ {page_num + 1} ===\n{page_text}\n"
                    
                    if text.strip():
                        result = {
                            'type': 'document_pdf',
                            'extracted_text': text.strip(),
                            'extraction_success': True,
                            'extraction_method': 'pdfplumber',
                            'pages': len(pdf.pages)
                        }
                        
                        # ë¬¸ì„œ ìš”ì•½ ìƒì„±
                        result['document_summary'] = self._summarize_document(
                            text, filename, 'PDF ë³´ê³ ì„œ'
                        )
                        
                        return result
                        
                except Exception as e:
                    print(f"[âš ï¸ pdfplumber ì‹¤íŒ¨] {str(e)}")
                finally:
                    try:
                        os.unlink(temp_file_path)
                    except:
                        pass
            
            # PyPDF2ë¡œ ì¬ì‹œë„
            if PYPDF2_AVAILABLE:
                return self._process_pdf_fallback(attachment_data, filename)
            
            return {'type': 'document_pdf', 'extraction_success': False, 'error': 'No PDF library available'}
            
        except Exception as e:
            print(f"[â—PDF ì²˜ë¦¬ ì˜¤ë¥˜] {str(e)}")
            return {'type': 'document_pdf', 'error': str(e), 'extraction_success': False}
    
    def _process_pdf_fallback(self, attachment_data, filename):
        """PDF ëŒ€ì²´ ì²˜ë¦¬ (PyPDF2)"""
        try:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
                temp_file.write(attachment_data)
                temp_file_path = temp_file.name
            
            with open(temp_file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n=== í˜ì´ì§€ {page_num + 1} ===\n{page_text}\n"
            
            if text.strip():
                result = {
                    'type': 'document_pdf',
                    'extracted_text': text.strip(),
                    'extraction_success': True,
                    'extraction_method': 'pypdf2',
                    'pages': len(pdf_reader.pages)
                }
                
                result['document_summary'] = self._summarize_document(
                    text, filename, 'PDF ë³´ê³ ì„œ'
                )
                
                return result
            
            return {'type': 'document_pdf', 'extraction_success': False}
            
        except Exception as e:
            return {'type': 'document_pdf', 'error': str(e), 'extraction_success': False}
        finally:
            try:
                os.unlink(temp_file_path)
            except:
                pass
    
    def _process_docx(self, attachment_data, filename):
        """Word ë¬¸ì„œ ì²˜ë¦¬"""
        if not DOCX_AVAILABLE:
            return {'type': 'document_word', 'error': 'python-docx not available', 'extraction_success': False}
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_file:
                temp_file.write(attachment_data)
                temp_file_path = temp_file.name
            
            doc = Document(temp_file_path)
            
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # í‘œ ë‚´ìš©ë„ ì¶”ì¶œ
            for table in doc.tables:
                text += "\n=== í‘œ ë°ì´í„° ===\n"
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            if text.strip():
                result = {
                    'type': 'document_word',
                    'extracted_text': text.strip(),
                    'extraction_success': True,
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables)
                }
                
                result['document_summary'] = self._summarize_document(
                    text, filename, 'Word ë¬¸ì„œ'
                )
                
                return result
            
            return {'type': 'document_word', 'extraction_success': False}
            
        except Exception as e:
            return {'type': 'document_word', 'error': str(e), 'extraction_success': False}
        finally:
            try:
                os.unlink(temp_file_path)
            except:
                pass
    
    def _process_pptx(self, attachment_data, filename):
        """PowerPoint ì²˜ë¦¬"""
        if not PPTX_AVAILABLE:
            return {'type': 'document_presentation', 'error': 'python-pptx not available', 'extraction_success': False}
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(attachment_data)
                temp_file_path = temp_file.name
            
            prs = Presentation(temp_file_path)
            
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n=== ìŠ¬ë¼ì´ë“œ {slide_num + 1} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    
                    if hasattr(shape, 'has_table') and shape.has_table:
                        text += "\n--- í‘œ ---\n"
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text += " | ".join(row_text) + "\n"
            
            if text.strip():
                result = {
                    'type': 'document_presentation',
                    'extracted_text': text.strip(),
                    'extraction_success': True,
                    'slides': len(prs.slides)
                }
                
                result['document_summary'] = self._summarize_document(
                    text, filename, 'PowerPoint í”„ë ˆì  í…Œì´ì…˜'
                )
                
                return result
            
            return {'type': 'document_presentation', 'extraction_success': False}
            
        except Exception as e:
            return {'type': 'document_presentation', 'error': str(e), 'extraction_success': False}
        finally:
            try:
                os.unlink(temp_file_path)
            except:
                pass
    
    def _process_xlsx(self, attachment_data, filename):
        """Excel ì²˜ë¦¬"""
        if not PANDAS_AVAILABLE:
            return {'type': 'document_spreadsheet', 'error': 'pandas not available', 'extraction_success': False}
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_file:
                temp_file.write(attachment_data)
                temp_file_path = temp_file.name
            
            xl_file = pd.ExcelFile(temp_file_path)
            
            text = ""
            total_rows = 0
            
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(temp_file_path, sheet_name=sheet_name)
                
                if not df.empty:
                    text += f"\n=== ì‹œíŠ¸: {sheet_name} ===\n"
                    text += "ì»¬ëŸ¼: " + " | ".join(str(col) for col in df.columns) + "\n\n"
                    
                    for idx, row in df.head(20).iterrows():
                        row_text = []
                        for value in row:
                            if pd.notna(value):
                                row_text.append(str(value))
                            else:
                                row_text.append("")
                        text += " | ".join(row_text) + "\n"
                    
                    total_rows += len(df)
                    
                    if len(df) > 20:
                        text += f"... (ì´ {len(df)}í–‰ ì¤‘ ì²˜ìŒ 20í–‰ë§Œ í‘œì‹œ)\n"
            
            if text.strip():
                result = {
                    'type': 'document_spreadsheet',
                    'extracted_text': text.strip(),
                    'extraction_success': True,
                    'sheets': len(xl_file.sheet_names),
                    'total_rows': total_rows
                }
                
                result['document_summary'] = self._summarize_document(
                    text, filename, 'Excel ìŠ¤í”„ë ˆë“œì‹œíŠ¸'
                )
                
                return result
            
            return {'type': 'document_spreadsheet', 'extraction_success': False}
            
        except Exception as e:
            return {'type': 'document_spreadsheet', 'error': str(e), 'extraction_success': False}
        finally:
            try:
                os.unlink(temp_file_path)
            except:
                pass
    
#     def _summarize_document(self, text, filename, file_type):
#         """ë¬¸ì„œ ìš”ì•½ ìƒì„± (Qwen 1.5-1.8B ëª¨ë¸ ì‚¬ìš©)"""
#         try:
#             if len(text) > 600:
#                 text = text[:600]
            
#             # Qwen ëª¨ë¸ë¡œ ìš”ì•½ ì‹œë„
#             if self.ai_models.load_qwen_model():
#                 try:
#                     prompt = f"""<|im_start|>system
# ë‹¹ì‹ ì€ ë¬¸ì„œ ìš”ì•½ ì „ë¬¸ê°€ì…ë‹ˆë‹¤.
# <|im_end|>
# <|im_start|>user
# ë‹¤ìŒì€ '{filename}' íŒŒì¼ì˜ ë‚´ìš©ì…ë‹ˆë‹¤. ì´ ë¬¸ì„œë¥¼ ìš”ì•½í•´ì£¼ì„¸ìš”.

# íŒŒì¼ í˜•ì‹: {file_type}
# ë‚´ìš©:
# {text}

# ìš”ì•½ ì§€ì¹¨:
# 1. ì£¼ìš” ë‚´ìš©ì„ 3-5ê°œ í¬ì¸íŠ¸ë¡œ ìš”ì•½
# 2. í•µì‹¬ í‚¤ì›Œë“œì™€ ìˆ˜ì¹˜ í¬í•¨
# 3. 150ì ì´ë‚´ë¡œ ê°„ê²°í•˜ê²Œ
# 4. í•œêµ­ì–´ë¡œ ì‘ë‹µ

# ìš”ì•½:
# <|im_end|>
# <|im_start|>assistant
# """
                    
#                     inputs = self.ai_models.qwen_tokenizer(prompt, return_tensors="pt").to(self.ai_models.qwen_model.device)
                    
#                     import torch
#                     with torch.no_grad():
#                         outputs = self.ai_models.qwen_model.generate(
#                             **inputs,
#                             max_new_tokens=150,
#                             temperature=0.3,
#                             do_sample=True,
#                             top_p=0.9,
#                             eos_token_id=self.ai_models.qwen_tokenizer.eos_token_id,
#                             pad_token_id=self.ai_models.qwen_tokenizer.pad_token_id
#                         )
                    
#                     generated_text = self.ai_models.qwen_tokenizer.decode(outputs[0], skip_special_tokens=True)
                    
#                     # "assistant" ì´í›„ í…ìŠ¤íŠ¸ë§Œ ê°€ì ¸ì˜¤ê¸°
#                     if "assistant" in generated_text:
#                         summary = generated_text.split("assistant")[-1].strip()
#                     else:
#                         summary = generated_text[len(prompt):].strip()
                    
#                     return summary if summary else text[:200] + "..."
                    
#                 except Exception as e:
#                     print(f"[âš ï¸ Qwen ë¬¸ì„œ ìš”ì•½ ì‹¤íŒ¨] {str(e)}")
            
#             # Qwen ì‹¤íŒ¨ ì‹œ ê°„ë‹¨í•œ ìš”ì•½ìœ¼ë¡œ fallback
#             sentences = text.split('.')
#             important_sentences = [s.strip() for s in sentences[:3] if len(s.strip()) > 10]
#             return '. '.join(important_sentences) + '.' if important_sentences else text[:200] + "..."
                
#         except Exception as e:
#             print(f"[âš ï¸ ë¬¸ì„œ ìš”ì•½ ì˜¤ë¥˜] {str(e)}")
#             return text[:200] + "..." if len(text) > 200 else text

    #0824 ìˆ˜ì •
     
    def _summarize_document(self, text, filename, file_type):
        """ë¬¸ì„œ ìš”ì•½ ìƒì„± (Qwen 1.5-1.8B ëª¨ë¸ ì‚¬ìš©)"""
        try:
        # 1) NPU(Genie) ê²½ë¡œ
            print("ë¬¸ì„œ ìš”ì•½ NPU ì„±ê³µ ")
            return genie_summarize_document(text, filename, file_type, max_words=25, max_chars=800)
        except Exception as ge:
            print(f"[âš ï¸ Genie ìš”ì•½ ì‹¤íŒ¨] {ge}")


        
        try:
            if len(text) > 600:
                text = text[:600]
            
            # Qwen ëª¨ë¸ë¡œ ìš”ì•½ ì‹œë„
            if self.ai_models.load_qwen_model():
                try:
                    prompt = f"""<|im_start|>system
ë‹¹ì‹ ì€ ë¬¸ì„œ ìš”ì•½ ì „ë¬¸ê°€ì…ë‹ˆë‹¤.
<|im_end|>
<|im_start|>user
ë‹¤ìŒì€ '{filename}' íŒŒì¼ì˜ ë‚´ìš©ì…ë‹ˆë‹¤. ì´ ë¬¸ì„œë¥¼ ìš”ì•½í•´ì£¼ì„¸ìš”.

íŒŒì¼ í˜•ì‹: {file_type}
ë‚´ìš©:
{text}

ìš”ì•½ ì§€ì¹¨:
1. ì£¼ìš” ë‚´ìš©ì„ 3-5ê°œ í¬ì¸íŠ¸ë¡œ ìš”ì•½
2. í•µì‹¬ í‚¤ì›Œë“œì™€ ìˆ˜ì¹˜ í¬í•¨
3. 150ì ì´ë‚´ë¡œ ê°„ê²°í•˜ê²Œ
4. í•œêµ­ì–´ë¡œ ì‘ë‹µ

ìš”ì•½:
<|im_end|>
<|im_start|>assistant
"""
                    
                    inputs = self.ai_models.qwen_tokenizer(prompt, return_tensors="pt").to(self.ai_models.qwen_model.device)
                    
                    import torch
                    with torch.no_grad():
                        outputs = self.ai_models.qwen_model.generate(
                            **inputs,
                            max_new_tokens=150,
                            temperature=0.3,
                            do_sample=True,
                            top_p=0.9,
                            eos_token_id=self.ai_models.qwen_tokenizer.eos_token_id,
                            pad_token_id=self.ai_models.qwen_tokenizer.pad_token_id
                        )
                    
                    generated_text = self.ai_models.qwen_tokenizer.decode(outputs[0], skip_special_tokens=True)
                    
                    # "assistant" ì´í›„ í…ìŠ¤íŠ¸ë§Œ ê°€ì ¸ì˜¤ê¸°
                    if "assistant" in generated_text:
                        summary = generated_text.split("assistant")[-1].strip()
                    else:
                        summary = generated_text[len(prompt):].strip()
                    
                    return summary if summary else text[:200] + "..."
                    
                except Exception as e:
                    print(f"[âš ï¸ Qwen ë¬¸ì„œ ìš”ì•½ ì‹¤íŒ¨] {str(e)}")
            
            # Qwen ì‹¤íŒ¨ ì‹œ ê°„ë‹¨í•œ ìš”ì•½ìœ¼ë¡œ fallback
            sentences = text.split('.')
            important_sentences = [s.strip() for s in sentences[:3] if len(s.strip()) > 10]
            return '. '.join(important_sentences) + '.' if important_sentences else text[:200] + "..."
                
        except Exception as e:
            print(f"[âš ï¸ ë¬¸ì„œ ìš”ì•½ ì˜¤ë¥˜] {str(e)}")
            return text[:200] + "..." if len(text) > 200 else text
      
    
    def _manage_cache_size(self):
        """ìºì‹œ í¬ê¸° ê´€ë¦¬"""
        if len(self.attachment_cache) > self.config.MAX_CACHE_SIZE:
            oldest_key = next(iter(self.attachment_cache))
            del self.attachment_cache[oldest_key]
            print(f"[ğŸ—‘ï¸ ìºì‹œ ì •ë¦¬] ì˜¤ë˜ëœ í•­ëª© ì‚­ì œ: {oldest_key}")
    
    def generate_attachment_summary(self, attachments):
        """ì²¨ë¶€íŒŒì¼ ìš”ì•½ ìƒì„±"""
        if not attachments:
            return ""
        
        total_files = len(attachments)
        
        # íŒŒì¼ íƒ€ì…ë³„ ë¶„ë¥˜
        images = [att for att in attachments if att.get('type') == 'image']
        documents = [att for att in attachments if att.get('type', '').startswith('document_')]
        others = [att for att in attachments if att.get('type') not in ['image'] and not att.get('type', '').startswith('document_')]
        
        summary_parts = []
        
        if images:
            total_objects = sum(att.get('object_count', 0) for att in images)
            ocr_texts = [att for att in images if att.get('ocr_success')]
            
            if total_objects > 0:
                summary_parts.append(f"ì´ë¯¸ì§€ {len(images)}ê°œ({total_objects}ê°œ ê°ì²´)")
            else:
                summary_parts.append(f"ì´ë¯¸ì§€ {len(images)}ê°œ")
                
            if ocr_texts:
                summary_parts.append(f"í…ìŠ¤íŠ¸ ì¶”ì¶œ {len(ocr_texts)}ê°œ")
        
        if documents:
            doc_types = {}
            successful_extractions = 0
            
            for doc in documents:
                doc_type = doc.get('type', '').replace('document_', '')
                doc_types[doc_type] = doc_types.get(doc_type, 0) + 1
                
                if doc.get('extraction_success'):
                    successful_extractions += 1
            
            for doc_type, count in doc_types.items():
                type_names = {
                    'pdf': 'PDF', 
                    'word': 'Word', 
                    'presentation': 'PPT', 
                    'spreadsheet': 'Excel'
                }
                type_name = type_names.get(doc_type, doc_type.upper())
                summary_parts.append(f"{type_name} {count}ê°œ")
            
            if successful_extractions > 0:
                summary_parts.append(f"ìš”ì•½ ê°€ëŠ¥ {successful_extractions}ê°œ")
        
        if others:
            summary_parts.append(f"ê¸°íƒ€ {len(others)}ê°œ")
        
        if summary_parts:
            return f"ğŸ“ {total_files}ê°œ íŒŒì¼: " + ", ".join(summary_parts)
        else:
            return f"ğŸ“ {total_files}ê°œ íŒŒì¼"
    
    def clear_cache(self):
        """ìºì‹œ ì´ˆê¸°í™”"""
        cache_count = len(self.attachment_cache)
        self.attachment_cache.clear()
        return cache_count
    
    def get_available_features(self):
        """ì‚¬ìš© ê°€ëŠ¥í•œ ê¸°ëŠ¥ ëª©ë¡ ë°˜í™˜"""
        return self.features